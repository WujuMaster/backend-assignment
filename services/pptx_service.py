from os import path

from deep_translator import GoogleTranslator
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.parts.image import Image
from pptx.text.text import TextFrame, _Paragraph
from pptx.util import Pt

from utils.helper import prepare_output_dir


class PPTXHandler:
    def __init__(
        self,
        data_dir: str,
        output_dir: str,
        src_lang: str = "en",
        target_lang: str = "vi",
    ):
        """Class for handling PPTX files - extracting images and translating text."""
        self.data_dir = data_dir
        self.output_dir = output_dir
        self.translator = GoogleTranslator(source=src_lang, target=target_lang)

    def translate_text(self, original_text: str):
        """Translate text using Google Translator API.

        Args:
            original_text: str, text to be translated

        Returns:
            str, translated text
        """
        return self.translator.translate(original_text)

    def add_subtext_to_paragraph(self, paragraph: _Paragraph, text: str):
        """Add a new run of subtext under the original paragraph run.

        Args:
            paragraph: _Paragraph object
            text: str, text to add under the paragraph
        """
        paragraph.line_spacing = 1.0
        for _ in range(paragraph.level):
            paragraph.space_before += Pt(8)  # type: ignore - library issue
        run = paragraph.add_run()
        run.text = f"\n{text}"
        run.font.size = Pt(8)

    def process_text_frame(self, text_frame: TextFrame):
        """Process a text frame by translating its text and adding the
        translated subtext under the original text.

        Args:
            text_frame: TextFrame object
        """
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        text_frame.word_wrap = True
        for paragraph in text_frame.paragraphs:
            original_text = paragraph.text
            processed_text = self.translate_text(original_text)
            self.add_subtext_to_paragraph(paragraph, processed_text)

    def extract_images(self, image: Image, slide_id: int, shape_id: int, img_dir: str):
        """Extract images from a shape and save them to the output directory.

        Args:
            image: Image object
            slide_id: int, slide ID
            shape_id: int, shape ID
            img_dir: str, output directory for saving images
        """
        image_bytes = image.blob
        image_filename = f"slide-{slide_id}_shape-{shape_id}"
        if image.filename:
            image_filename += f"_{image.filename.split('.')[0]}"
        img_ext = image.ext or "png"
        output_file = path.join(img_dir, f"{image_filename}.{img_ext}")
        with open(output_file, "wb") as f:
            f.write(image_bytes)

    def extract_images_and_translate_text(self, pptx_file: str):
        """Extract images and translate text of a PPTX file.

        Args:
            pptx_file: str, path to the PPTX file
        """
        img_dir = prepare_output_dir(path.basename(pptx_file), self.output_dir)

        prs = Presentation(pptx_file)
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:  # type: ignore - library issue
                    self.extract_images(
                        shape.image,
                        slide.slide_id,
                        shape.shape_id,
                        img_dir=img_dir,
                    )
                if shape.has_text_frame:
                    self.process_text_frame(shape.text_frame)
        prs.save(pptx_file)
